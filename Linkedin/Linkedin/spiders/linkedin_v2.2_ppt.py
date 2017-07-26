import re

import optparse
import MySQLdb
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor
import urllib
import os
from pptx.dml.fill import FillFormat
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from ppt_db_constants import *
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN

class Login(object):

    def __init__(self, options):
        self.con = MySQLdb.connect(db=db_name,
                      user=db_user,
                      passwd=db_passwd,
                      charset="utf8",
                      host=db_host,
                      use_unicode=True)
        self.cur = self.con.cursor()
        self.prs = Presentation()
        self.prs.slide_width = Cm(33.87)
        self.prs.slide_height = Cm(19.05)
        self.select_qry = select_qry
        self.select_qry1 = select_qry1_v2
        self.select_qry2 = select_qry2
        self.select_qry3 = select_qry3_v2_1
        self.select_qry4 = select_qry4
        self.select_qry5 = select_qry5_v2
        self.select_qry6 = select_qry6_v2
        self.main()

    def main(self):
        count = 0
        member = options.memberid.split(',')
        no_image = ''
        for member_id in member:
                self.cur.execute(self.select_qry % member_id)
		rows = self.cur.fetchall()
                if not rows: 
                    print "Invalid member_id %s .Please try by giving valid member_id " % member_id
                    continue
                else:
                    rows = rows[0]
		    name, sk, summary, headline, location, profile_url = rows
		    self.cur.execute(self.select_qry2 % member_id)
		    profile_image = self.cur.fetchall()
		    try: 
                        path = profile_image[0][1]
                        profile_image = profile_image[0][0]
		    except: 
                        profile_image = ''
                        no_image = 'no image'
		    self.get_ppt(rows, profile_image, sk, member_id, path, no_image)

    def get_ppt(self, records, profile_image, sk, member_id, path, no_image):
        self.cur.execute(self.select_qry1 % sk)
        exp_data = self.cur.fetchall()
        #Selecting the blank layout(6)
	title_only_slide_layout = self.prs.slide_layouts[6]
	slide = self.prs.slides.add_slide(title_only_slide_layout)
	shapes = slide.shapes
        import pdb;pdb.set_trace()
        self.add_title(slide, shapes, records[0], exp_data[0][0], sk, records[4], exp_data)
        self.add_summary(slide, shapes, records[2], records[5])
        img_path = self.download_image(profile_image, member_id, path)
        self.add_tables(slide, shapes, exp_data)
        self.add_languages(slide, shapes, sk)
        self.add_education(slide, shapes, sk)
        self.prof_exp_table(shapes, slide, exp_data)
        self.copyrights(slide, shapes) 
        self.add_connector(slide, shapes, Cm(0.00), Cm(17.90), Cm(33.90), Cm(0.0))
        self.insert_image(slide, shapes, img_path, member_id, no_image)
        self.prs.save(options.filename)
        print "ppt succesfully created for %s" % member_id

    def mergeCellsVertically(self, table, start_row_idx, end_row_idx, col_idx):
        row_count = end_row_idx - start_row_idx + 1
        column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]
        column_cells[0]._tc.set('rowSpan', str(row_count))
        for c in column_cells[1:]:
            c._tc.set('vMerge', '1')

    def merge_cells(self, row_idx, start_col_idx, end_col_idx, table2):
        col_count = end_col_idx - start_col_idx + 1
        row_cells = [c for c in table2.rows[row_idx].cells][start_col_idx:end_col_idx]
        row_cells[0]._tc.set('gridSpan', str(col_count))
        for c in row_cells[1:]:
            c._tc.set('hMerge', '2')


    def add_tables(self, slide, shapes, exp_data):
        txBox = slide.shapes.add_textbox(left=Inches(2.3), top=Cm(2.44), width=Inches(4), height=Cm(1.2))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if exp_data[0][2]: 
            tf.text = exp_data[0][2].split('-')[0]+'-current'
            tf.word_wrap = True
	    tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
	    tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.italic = True
        try: self.add_current_com_logo(slide, shapes, exp_data[0][0])
        except: print  "No company logo Found"
        return slide, shapes

    def add_connector(self, slide, shapes, left, top, width, height):
        line = shapes.add_shape(MSO_CONNECTOR.STRAIGHT, left, top, width, height)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(191, 191, 191)
        return slide, shapes

    def add_current_com_logo(self, slide, shapes, company_name):
        self.cur.execute(select_qry4_v2 % company_name)
        logo = self.cur.fetchall()
        image = urllib.URLopener()
        try: 
              image_name = image.retrieve(str(logo[0][0]), 'logo.jpg')
              img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
              pic = slide.shapes.add_picture(img_path,\
              left=Cm(13.93), top=Cm(0.10), width=Cm(1.30), height=Cm(1.30))
        except: print "Company logo not Found"

        return slide, shapes

    def add_current_role(self, slide, shapes, exp_data):
        table2 = shapes.add_table(rows=2, cols=2, \
        left=Cm(14), top=Cm(6.30), width=Inches(2.3),\
        height=Inches(1.0)).table
        row_idx, start_col_idx, end_col_idx = 0, 0, 1
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table2)
        # column widths
        table2.columns[0].width = Inches(2.4)
        table2.columns[1].width = Inches(1.2)
        table2.rows[0].height = Inches(0.3)
        table2.rows[1].height = Inches(0.3)
        if exp_data[0][0]: table2.cell(0, 0).text = 'Current role' + ' at '+ exp_data[0][0]
        else: table2.cell(0, 0).text = "Current role"
        table2.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(13)
        table2.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        table2.cell(1, 0).text = exp_data[0][1]
        table2.cell(1, 0).text_frame.paragraphs[0].font.size = Pt(9)
        table2.cell(1, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        if exp_data[0][2]: table2.cell(1, 1).text = exp_data[0][2].split('-')[0]+'-current'
        table2.cell(1, 1).text_frame.paragraphs[0].font.size = Pt(9)
        return slide, shapes

    def add_languages(self, slide, shapes, sk):
        self.cur.execute(self.select_qry6 % sk)
        lan = self.cur.fetchall()
        lan = lan[0][0]
        table6 = shapes.add_table(rows=2, cols=1, \
        left=Cm(25.50), top=Cm(12), width=Inches(2.3),\
        height=Inches(1.0)).table
        table6.columns[0].width = Inches(3.2)
        table6.cell(0, 0).text = 'Languages known'
        table6.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(11)
        table6.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        if options.languages == 'yes':
	        if lan: 
		    lan = lan.split(',')
		    if len(lan) >= 6:
			table6.cell(1, 0).text = lan[0] +'    ' + lan[1] + '\n\n\n' + \
						 lan[2]+'    '+lan[3] +  '\n\n\n' + \
						 lan[4]+'    '+lan[5] 
		    if len(lan) == 5:
			table6.cell(1, 0).text = lan[0] +'    ' + lan[1] + '\n\n\n' + \
						 lan[2]+'    '+lan[3] +  '\n\n\n' + \
						 lan[4]
		    if len(lan) == 4:
			 table6.cell(1, 0).text = lan[0] +'    ' + lan[1] + '\n\n\n' + \
						 lan[2]+'    '+lan[3] +  '\n\n\n' 
						 
		    if len(lan) == 3: 
			 table6.cell(1, 0).text = lan[0] +'    ' + lan[1] + '\n\n\n' + \
						 lan[2]

		    if len(lan) == 2:
			  table6.cell(1, 0).text = lan[0] +'    ' + lan[1] 
			    
		    if len(lan) == 1:
			  table6.cell(1, 0).text = lan[0] 
						 
                    table6.cell(1, 0).text_frame.paragraphs[0].font.size = Pt(9)
        image = urllib.URLopener()
        lan_icon =  'http://www.languageicon.org/LanguageIcon.jpg'
        image_name = image.retrieve(lan_icon, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        logo = slide.shapes.add_picture(img_path,
              left=Cm(32.50), top=Cm(12), width=Cm(0.80), height=Cm(0.80))

        return slide, shapes 

    def add_education(self, slide, shapes, sk):
        self.cur.execute(self.select_qry3 % sk)
        edu_data = self.cur.fetchall()
        table3 = shapes.add_table(rows=5, cols=1, \
        left=Cm(15.80), top=Cm(12), width=Inches(2.3),\
        height=Inches(1.0)).table
        row_idx, start_col_idx, end_col_idx = 0, 0, 1
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table3)
        table3.columns[0].width = Inches(3.7)
        table3.rows[0].height = Inches(0.3)
        table3.rows[1].height = Inches(0.4)
        table3.rows[2].height = Inches(0.4)
        table3.rows[3].height = Inches(0.4)
        table3.rows[4].height = Inches(0.4)
        table3.cell(0, 0).text = "Education"
        table3.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(13)
        table3.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        if options.education == 'yes':
		for ind, colu in enumerate(edu_data):
			row = 1 
			for inded, i in enumerate(edu_data):
				if ind == 0:
				    table3.cell(row, 0).text = i[0]+','+i[1]+','+i[2]
				    table3.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(9)
				    row = row+1
        image = urllib.URLopener() 
        edu_icon = 'https://tse4.mm.bing.net/th?id=OIP.jNa_Rgg0ydpGcxK8MxhiaAEsEs&pid=15.1&P=0&w=100&h=100'
        image_name = image.retrieve(edu_icon, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        logo = slide.shapes.add_picture(img_path,
              left=Cm(24), top=Cm(12), width=Cm(0.80), height=Cm(0.80))

        return slide, shapes

    def add_honors(self, slide, shapes, sk): 
        table4 = shapes.add_table(rows=3, cols=2, \
         left=Cm(14), top=Cm(13), width=Inches(2.3),\
        height=Inches(1.0)).table
        row_idx, start_col_idx, end_col_idx = 0, 0, 1
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table4)
        table4.columns[0].width = Inches(2.4)
        table4.columns[1].width = Inches(1.2)
        table4.rows[0].height = Inches(0.3)
        table4.rows[1].height = Inches(0.3)
        table4.cell(0, 0).text = "Honors & Awards"
        table4.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(13)
        table4.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        table4.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        table4.cell(0, 1).fill.background()
        self.cur.execute(self.select_qry5 % sk)
        honor_data = self.cur.fetchall()
        for ind, colu in enumerate(honor_data):
                row = 1        
                for inded, i in enumerate(honor_data):
                        if ind == 0:
                            table4.cell(row, 0).text = i[0]
                            table4.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(9)
                            table4.cell(row, 1).text = i[1]
                            table4.cell(row, 1).text_frame.paragraphs[0].font.size = Pt(9)
                            row += 1
        return slide, shapes

    def prof_exp_table(self, shapes, slide, exp_data): 
        txBox = slide.shapes.add_textbox(left=Cm(15.80), top=Cm(0.50), width=Inches(4), height=Cm(1))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'Prior Experiences'
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        table5 = shapes.add_table(rows=8, cols=4, \
         left=Cm(15.80), top=Cm(1.5), width=Cm(16.34),\
        height=Cm(13)).table
        table5.columns[0].width = Inches(1.3)
        table5.columns[1].width = Inches(2)
        table5.columns[2].width = Inches(2.3)
        table5.columns[3].width = Inches(1.5)
        table5.rows[0].height = Inches(0.2)
        table5.rows[1].height = Inches(0.5)
        table5.rows[2].height = Inches(0.5)
        table5.rows[3].height = Inches(0.5)
        table5.rows[4].height = Inches(0.5)
        table5.rows[5].height = Inches(0.5)
        table5.rows[6].height = Inches(0.5)
        table5.rows[7].height = Inches(0.5)
        table5.cell(0, 0).text = "Date"
        table5.cell(0, 1).text = 'company'
        table5.cell(0, 2).text = 'Role'
        table5.cell(0, 3).text = 'location'
        table5.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(9)
        table5.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(9)
        table5.cell(0, 2).text_frame.paragraphs[0].font.size = Pt(9)
        table5.cell(0, 3).text_frame.paragraphs[0].font.size = Pt(9)
        self.add_connector(slide, shapes, Cm(15.83), Cm(1.20), Cm(10), Cm(0.0))
        if options.experiences == 'yes': 
		for ind, colu in enumerate(exp_data):
			row = 1
		     
			for indk, i in enumerate(exp_data):
			    if ind == 0:
				table5.cell(row, 0).text = str(i[2].split('-')[0]+'-'+i[3].split('-')[0]).strip('-')
				table5.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(9)
				row = row+1
			    if ind == 1:
				table5.cell(row, 1).text = i[0]
				table5.cell(row, 1).text_frame.paragraphs[0].font.size = Pt(9)
				row = row+1
			    if ind == 2:
				table5.cell(row, 2).text = i[1]
				table5.cell(row, 2).text_frame.paragraphs[0].font.size = Pt(9)
				row = row+1
			    if ind == 3:
				table5.cell(row, 3).text = i[4]
				table5.cell(row, 3).text_frame.paragraphs[0].font.size = Pt(9)
				row = row+1
        image = urllib.URLopener()
        prof_icon = 'https://upload.wikimedia.org/wikipedia/commons/thumb/b/bb/Summary_icon.svg/2000px-Summary_icon.svg.png'
        image_name = image.retrieve(prof_icon, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        logo = slide.shapes.add_picture(img_path,
              left=Cm(24), top=Cm(0.50), width=Cm(0.80), height=Cm(0.60))
        return slide, shapes

    def download_image(self, profile_image, member_id, path):
        image = urllib.URLopener()
        pattern = ''
        if '/mpr/mpr/shrink' in profile_image: 
            pattern = "".join(re.findall('mpr/mpr/shrink_\d+_\d+', profile_image))
        if options.yes == 'yes': 
                if pattern: profile_image = profile_image.replace(pattern, 'media')
                else: profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                try: 
                    image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                    img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
                except: 
                    print "404 status found for %s while displaying image" % member_id
                    profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                    image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                    img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        if options.yes == 'no' and path:
                img_path = path
        if options.yes == 'no' and not path:
                profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        return img_path
        

    def insert_image(self, slide, shapes, img_path, member_id, no_image):
        #Inserting persons  image
        if options.silhoutte == 'yes': 
            sil = 'https://image.freepik.com/free-icon/male-user-shadow_318-34042.png'
            image = urllib.URLopener()
            image_name = image.retrieve(sil, 'sil_logo.jpg')
            img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        else: img_path = img_path
        pic2 = slide.shapes.add_picture(img_path, \
        left=Cm(0), top=Cm(0), width=Cm(5.20), height=Cm(5.20))
        cal = 'http://icons.iconarchive.com/icons/paomedia/small-n-flat/1024/calendar-icon.png'
        logo = 'https://positivemoves.com/wp-content/uploads/2017/03/Pmoves-1.png'
        location_icon = 'http://imageog.flaticon.com/icons/png/512/61/61121.png?size=50x50f&pad=10,10,10,10&ext=png&bg=FFFFFFFF'
        image = urllib.URLopener()
        image_name = image.retrieve(location_icon, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        pic = slide.shapes.add_picture(img_path,\
              left=Inches(2.3), top=Cm(4.20), width=Cm(1.20), height=Cm(1.20))
        cal_image = image.retrieve(cal, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(cal_image[0]))+ '/' +cal_image[0]
        pic3 = slide.shapes.add_picture(img_path,\
              left=Inches(2.3), top=Cm(3.40), width=Cm(1.20), height=Cm(0.97))
        logo = 'https://positivemoves.com/wp-content/uploads/2017/03/Pmoves-1.png'
        image = urllib.URLopener()
        image_name = image.retrieve(logo, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        logo = slide.shapes.add_picture(img_path, \
              left=Inches(11.81), top=Cm(18.20), width=Cm(3.20), height=Cm(0.85))
        return slide, shapes

    def add_summary(self, slide, shapes, summary, profile_url):
        self.add_connector(slide, shapes, Cm(0.30), Cm(6.45), Cm(12), Cm(0.0))
        txBox_sum = slide.shapes.add_textbox(left=Inches(0.2), top=Cm(5.40), \
        width=Inches(2), height=Inches(0.5))
        tf = txBox_sum.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'Summary'
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        tf.paragraphs[0].font.bold = True
        if options.summary == 'yes':
		txBox_ = slide.shapes.add_textbox(left=Cm(0), \
		top=Cm(6.5), width=Cm(14), height=Cm(5))
		tf = txBox_.text_frame
		tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
		tf.text = summary
		tf.word_wrap = True
		tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
		tf.paragraphs[0].font.size = Pt(9)
        image = urllib.URLopener()
        icon = 'https://tse3.mm.bing.net/th?id=OIP.YEIpYQpZ3NHiNizAZNWJVQDzEB&pid=15.1&P=0&w=300&h=300'
        image_name = image.retrieve(icon, 'keylogo.jpg')
        img = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        logo = slide.shapes.add_picture(img, left=Cm(11.40), top=Cm(5.50), width=Cm(0.80), height=Cm(0.80))

    def add_shapes(self, slide, shapes, RECTANGLE, left, top, width, height):
        shape1 = slide.shapes.add_shape(RECTANGLE, left, top,\
        width, height)
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(216, 191, 216)
        return slide, shapes

    def add_title(self, slide, shapes, title, org, sk, location, exp_data):
        #Adding Title
        txBox = slide.shapes.add_textbox(left=Inches(2.3), top=Cm(0.10), width=Inches(4), height=Cm(1))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if options.name == 'yes':
            tf.text = title
        else : tf.text = 'XXXXXX'
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(18.2)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        txBox2 = slide.shapes.add_textbox(left=Inches(2.3), top=Cm(1), width=Inches(3.5), height=Cm(1.9))
        tf = txBox2.text_frame
        tf.text = exp_data[0][1] + ' at ' + exp_data[0][0]
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(14)
        txBox3 = slide.shapes.add_textbox(left=Inches(2.67), top=Cm(3.60), width=Inches(1.5), height=Cm(1))
        tf = txBox3.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        exp_dur = self.exp_duration(slide, shapes, sk)
        tf.text = exp_dur
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        txBox4 = slide.shapes.add_textbox(left=Inches(2.67), top=Cm(4.60), width=Inches(2.7), height=Cm(1))
        tf = txBox4.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = location
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        return slide, shapes

    def exp_duration(self, slide, shapes, sk):
        self.cur.execute(self.select_qry4 % sk)
        exp_dur = self.cur.fetchall()
        exp_list = []
        if exp_dur:
            for exp in exp_dur:
               if 'years' in exp[0]:
                   dura = re.findall('\d+', exp[0])[0]
                   exp_list.append(int(dura))
               else: exp_duration = 'Below 2 years'
            exp_duration = str(sum(exp_list))+'+ Years'
            if '0' in exp_duration: exp_duration = 'Experience N/A'
            return exp_duration
    
    def copyrights(self, slide, shapes):
        txtBox = slide.shapes.add_textbox(left=Cm(2.9), top=Inches(6.9), width=Cm(15), height=Cm(1.2))
        tf = txtBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
	run.text = u'\xa9 2017 Positive Moves Consulting | Web: www.positivemoves.com'
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(12)
        font.color.rgb = RGBColor(0, 0, 0)
 

    def __del__(self):
        self.con.close()
        self.cur.close()

if __name__ == '__main__':
    parser = optparse.OptionParser()
    parser.add_option('-m', '--memberid', default=None, help='member_id, one or many separated by commas')
    parser.add_option('-f', '--filename', default='linkedin_member.pptx', help='filename, give any filename')
    parser.add_option('-i', '--yes', default='no', help='high resolution image /normal image')
    parser.add_option('-t', '--silhoutte', default='no', help='high resolution image /normal image')
    parser.add_option('-s', '--summary', default='no', help='give summary if you want')
    parser.add_option('-e', '--education', default='no', help='give education if you want')
    parser.add_option('-p', '--experiences', default='no', help='give experiences if you want')
    parser.add_option('-l', '--languages', default='no', help='give language if you want')
    parser.add_option('-n', '--name', default='no', help='give name if you want to display name ')
    (options, args) = parser.parse_args()
    Login(options)

