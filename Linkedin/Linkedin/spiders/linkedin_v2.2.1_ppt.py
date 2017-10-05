#version  linkedin_v2.2_ppt.py v2.2.1 2017/21/09 5:30:30 
#Modified the font color, coordinates according to gridlines view.
#line numbers :131,140,170,238,324,348,351,353,469,520,
#449-460,306-316,350-353,158-171,438-441 

import re
import optparse
import MySQLdb
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor
import urllib
import os
from pptx.dml.fill import FillFormat
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from constants import *
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN

class Login(object):

    def __init__(self, options):
        self.con = MySQLdb.connect(db=db_name,
                      user=db_user,
                      passwd=db_passwd,
                      charset="utf8",
                      host=db_host,
                      use_unicode=True)
        self.cur = self.con.cursor()
        self.prs = Presentation()
        self.prs.slide_width = Cm(33.87)
        self.prs.slide_height = Cm(19.05)
        self.select_qry_member = select_qry_member
        self.select_qry = select_qry
        self.select_exp_qry = select_exp_qry
        self.select_edu_qry = select_edu_qry
        self.select_dur_qry = select_dur_qry         
        self.select_qry4_v2 = select_qry4_v2
        self.select_qry6 = select_qry6_v2
        self.select_qry7 = select_qry7_v2
        self.main()

    def main(self):
        count = 0
        member = []
        member_id = self.cur.execute(self.select_qry_member)
        member_id = self.cur.fetchall() 
        page_num = 0
        for i in member_id : 
            member.append(i[0])
        if options.memberid == 'no' : member = member
        else : member = options.memberid.split(',')
        no_image = ''
        for member_id in member:
                page_num += 1
                self.cur.execute(self.select_qry % member_id)
		rows = self.cur.fetchall()
                if not rows: 
                    print "Invalid member_id %s .Please try by giving valid member_id " % member_id
                    continue
                else:
                    rows = rows[0]
		    name, summary, headline, location, profile_url, profile_image = rows
		    if profile_image : 
                        path = ''
                        profile_image = profile_image
		    else: 
                        profile_image = ''
                        path = ''
		    self.get_ppt(rows, profile_image, member_id, path, page_num)


    def get_ppt(self, records, profile_image, member_id, path, page_num):
        self.cur.execute(self.select_exp_qry % member_id)
        exp_data = self.cur.fetchall()
        #Selecting the blank layout(6)
	title_only_slide_layout = self.prs.slide_layouts[6]
	slide = self.prs.slides.add_slide(title_only_slide_layout)
	shapes = slide.shapes
        self.add_title(slide, shapes, records[0], exp_data[0][0], records[3], exp_data, member_id)
        if options.summary == 'yes': self.add_summary(slide, shapes, records[1], records[5])
        img_path = self.download_image(profile_image, member_id, path)
        self.add_tables(slide, shapes, exp_data, member_id)
        if options.languages == 'yes': self.add_languages(slide, shapes, member_id)
        if options.education == 'yes': self.add_education(slide, shapes, member_id)
        if options.experiences == 'yes': self.prof_exp_table(shapes, slide, exp_data)
        self.copyrights(slide, shapes, page_num) 
        self.add_connector(slide, shapes, Cm(0.00), Cm(17.90), Cm(33.90), Cm(0.0))
        self.insert_image(slide, shapes, img_path, member_id)
        self.prs.save(options.filename)
        print "ppt succesfully created for %s" % member_id

    def mergeCellsVertically(self, table, start_row_idx, end_row_idx, col_idx):
        row_count = end_row_idx - start_row_idx + 1
        column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]
        column_cells[0]._tc.set('rowSpan', str(row_count))
        for c in column_cells[1:]:
            c._tc.set('vMerge', '1')

    def merge_cells(self, row_idx, start_col_idx, end_col_idx, table2):
        col_count = end_col_idx - start_col_idx + 1
        row_cells = [c for c in table2.rows[row_idx].cells][start_col_idx:end_col_idx]
        row_cells[0]._tc.set('gridSpan', str(col_count))
        for c in row_cells[1:]:
            c._tc.set('hMerge', '2')


    def add_tables(self, slide, shapes, exp_data, member_id):
        headline = ''
        if  exp_data[0][1] and  exp_data[0][0] : headline = exp_data[0][1] + ' at ' + exp_data[0][0] 
        if len(headline) > 95 :  
            txBox = slide.shapes.add_textbox(left=Cm(5.80), top=Cm(3.20), width=Inches(4), height=Cm(1.2))
        else : txBox = slide.shapes.add_textbox(left=Cm(5.80), top=Cm(3.10), width=Inches(4), height=Cm(1.2))

        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if exp_data[0][2]: 
            tf.text = exp_data[0][2].split('-')[0]+'-current'
        else : tf.text = 'current year N/A'
        tf.word_wrap = True
	tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
	tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105) 
        tf.paragraphs[0].font.italic = True
 
        try: self.add_current_com_logo(slide, shapes, member_id)
        except: print  "No company logo Found"
        return slide, shapes

    def add_connector(self, slide, shapes, left, top, width, height):
        line = shapes.add_shape(MSO_CONNECTOR.STRAIGHT, left, top, width, height)
        line.line.color.rgb = RGBColor(105, 105, 105)
        return slide, shapes

    def add_current_com_logo(self, slide, shapes, member_id):
        self.cur.execute(self.select_qry4_v2 % member_id)
        logo = self.cur.fetchall()
        image = urllib.URLopener()
        try: 
              image_name = image.retrieve(str(logo[0][0]), 'logo.jpg')
              img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
              pic = slide.shapes.add_picture(img_path,\
              left=Cm(10.60), top=Cm(3.20), width=Cm(1.65), height=Cm(1.65))
        except: print "Company logo not Found"

        return slide, shapes

    def add_education(self, slide, shapes, member_id):
        self.cur.execute(self.select_edu_qry % member_id)
        edu_data = self.cur.fetchall()
        table3 = shapes.add_table(rows=5, cols=1, \
        left=Cm(16.20), top=Cm(12.95), width=Inches(2.3),\
        height=Inches(1.0)).table
        row_idx, start_col_idx, end_col_idx = 0, 0, 1
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table3)
        table3.columns[0].width = Inches(3.7)
        table3.rows[0].height = Inches(0.3)
        table3.rows[1].height = Inches(0.4)
        table3.rows[2].height = Inches(0.4)
        table3.rows[3].height = Inches(0.4)
        table3.rows[4].height = Inches(0.4)
        table3.cell(0, 0).text = "ducation"
        table3.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(16)
        table3.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(105 ,105, 105)
        table3.cell(0, 0).text_frame.paragraphs[0].font.bold = False
        table3.cell(0, 0).text_frame.paragraphs[0].font.name = 'Gadugi'
        txBox_ = slide.shapes.add_textbox(left=Cm(15.79), top=Cm(12.80), width=Cm(1.60), height=Cm(1.75))
        tf = txBox_.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'E'
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        row = 1
        try : 
            edu_data = edu_data[0]
            edu_data = edu_data[0:3],edu_data[3:6],edu_data[6:9]
        except : print "no_education data present in our Database"
 
        if edu_data : 
	    for i in edu_data : 
                if i :
                    table3.cell(row, 0).text = i[0]+','+i[1]+','+i[2]
                    table3.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(9)
                    table3.cell(row, 0).text_frame.paragraphs[0].font.name = 'Gadugi'
	            row = row+1

        for row in table3.rows:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        self.add_connector(slide, shapes, Cm(16.00), Cm(13.80), Cm(7.07), Cm(0.00)) 
        edu_icon = os.getcwd()+'/static_images/'+'education.jpg'
        logo = slide.shapes.add_picture(edu_icon,
              left=Cm(21.53), top=Cm(12.80), width=Cm(1.12), height=Cm(1.00))

        return slide, shapes

    def add_honors(self, slide, shapes, sk): 
        table4 = shapes.add_table(rows=3, cols=2, \
         left=Cm(14), top=Cm(13), width=Inches(2.3),\
        height=Inches(1.0)).table
        row_idx, start_col_idx, end_col_idx = 0, 0, 1
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table4)
        table4.columns[0].width = Inches(2.4)
        table4.columns[1].width = Inches(1.2)
        table4.rows[0].height = Inches(0.3)
        table4.rows[1].height = Inches(0.3)
        table4.cell(0, 0).text = "Honors & Awards"
        table4.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(13)
        table4.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        table4.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        table4.cell(0, 1).fill.background()
        self.cur.execute(self.select_qry5 % sk)
        honor_data = self.cur.fetchall()
        for ind, colu in enumerate(honor_data):
                row = 1        
                for inded, i in enumerate(honor_data):
                        if ind == 0:
                            table4.cell(row, 0).text = i[0]
                            table4.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(9)
                            table4.cell(row, 1).text = i[1]
                            table4.cell(row, 1).text_frame.paragraphs[0].font.size = Pt(9)
                            row += 1
        return slide, shapes

    def prof_exp_table(self, shapes, slide, exp_data_): 
        txBox = slide.shapes.add_textbox(left=Cm(15.52), top=Cm(0.06), width=Cm(1.27), height=Cm(1))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'P'
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        txBox_ = slide.shapes.add_textbox(left=Cm(15.94), top=Cm(0.20), width=Cm(10.16), height=Cm(1))
        tf = txBox_.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'rior Experiences'
        tf.word_wrap = True
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
        table5 = shapes.add_table(rows=8, cols=4, \
         left=Cm(15.80), top=Cm(1.64), width=Cm(16.34),\
        height=Cm(10.5)).table
        table5.columns[0].width = Inches(1.3)
        table5.columns[1].width = Inches(2)
        table5.columns[2].width = Inches(2.3)
        table5.columns[3].width = Inches(1.5)
        table5.rows[0].height = Inches(0.3)
        table5.rows[1].height = Inches(0.3)
        table5.rows[2].height = Inches(0.3)
        table5.rows[3].height = Inches(0.3)
        table5.rows[4].height = Inches(0.3)
        table5.rows[5].height = Inches(0.3)
        table5.rows[6].height = Inches(0.3)
        table5.rows[7].height = Inches(0.3)
        table5.cell(0, 0).text = "Date"
        table5.cell(0, 1).text = 'Company'
        table5.cell(0, 2).text = 'Role'
        table5.cell(0, 3).text = 'Location'
        list_ = [table5.cell(0, 0),table5.cell(0, 1),table5.cell(0, 2),table5.cell(0, 3)]
        for i in list_ :
            i.text_frame.paragraphs[0].font.size = Pt(10.5)
            i.text_frame.paragraphs[0].font.name = 'Gadugi'
            i.text_frame.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
        self.add_connector(slide, shapes, Cm(15.80), Cm(1.08), Cm(7.38), Cm(0))
        line = shapes.add_shape(MSO_CONNECTOR.STRAIGHT, left=Cm(15.80), top=Cm(1.45), width=Cm(18), height=Cm(0))
        line.line.color.rgb = RGBColor(105, 105, 105)
        line1 = shapes.add_shape(MSO_CONNECTOR.STRAIGHT, left=Cm(15.80), top=Cm(2.27), width=Cm(18), height=Cm(0))
        line1.line.color.rgb = RGBColor(105, 105, 105)
        row = 1
        try : exp_data = exp_data_[0]
        except : print  exp_data_
        exp_data = exp_data[0:5],exp_data[5:10],exp_data[10:15],\
        exp_data[15:20],exp_data[20:25],exp_data[25:30],\
        exp_data[30:35] 
        for i in exp_data :
            table5.cell(row, 0).text = str(i[2].split('-')[0]+'-'+i[3].split('-')[0]).strip('-')
            table5.cell(row, 0).text_frame.paragraphs[0].font.size = Pt(10)
            table5.cell(row, 0).text_frame.paragraphs[0].font.name = 'Gadugi'
            table5.cell(row, 1).text = i[0]
            table5.cell(row, 1).text_frame.paragraphs[0].font.size = Pt(10)
            table5.cell(row, 1).text_frame.paragraphs[0].font.name = 'Gadugi'
            table5.cell(row, 2).text = i[1][0:60]
            table5.cell(row, 2).text_frame.paragraphs[0].font.size = Pt(10)
            table5.cell(row, 2).text_frame.paragraphs[0].font.name = 'Gadugi'
            table5.cell(row, 3).text = i[4]
            table5.cell(row, 3).text_frame.paragraphs[0].font.size = Pt(10)
            table5.cell(row, 3).text_frame.paragraphs[0].font.name = 'Gadugi'
            row = row+1
        for row in table5.rows:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        prof_icon = os.getcwd()+'/static_images/'+'experience.jpg'
        logo = slide.shapes.add_picture(prof_icon,
              left=Cm(21.85), top=Cm(0.25), width=Cm(0.60), height=Cm(0.60))
       
        path =  os.path.dirname
        return slide, shapes

    def fb_gender(self, fb_gender):
        if fb_gender == 'male' :
            img_path =  os.getcwd()+'/static_images/'+ 'male.jpg'
        elif fb_gender == 'female':
            img_path =  os.getcwd()+'/static_images/'+'female.jpg'
        else: img_path =  os.getcwd()+'/static_images/'+'silhouette.jpeg'
        return img_path

    def download_image(self, profile_image, member_id, path):
        # Displaying Silhouette image based on gender
        try :
            self.cur.execute(self.select_qry7 % member_id)
            value = self.cur.fetchall()
            fb_gender = value[0][0]
            fb_gender = fb_gender.lower()
        except : fb_gender = ''
        # If silhoutee option is yes
        if options.silhoutte == 'yes': img_path = self.fb_gender(fb_gender)
        if options.silhoutte == 'yes' and  options.yes =='yes':
            print "silhouette and  -i yes options should not both be true"
            pass
        #If silhoutte option is no and need high resolution image
        image = urllib.URLopener()  
        pattern = ''
        if options.silhoutte == 'no' and options.yes == 'yes':
            if '/mpr/mpr/shrink' in profile_image: pattern = "".join(re.findall('mpr/mpr/shrink_\d+_\d+', profile_image))
            if pattern:
                    try :
                        profile_image = profile_image.replace(pattern, 'media')
                        image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
                    except:
                        print "404 status found for %s while displaying image" % member_id
                        # this line is for if pattern is there but that image is not opening and getting 404 error
                        img_path = self.fb_gender(fb_gender)

            else: img_path = self.fb_gender(fb_gender)

        #If silhoutte option is no and need normal image/no image path 
        elif options.silhoutte == 'no' and options.yes == 'no':
                img_path = self.fb_gender(fb_gender)
        return img_path
        
    def insert_image(self, slide, shapes, img_path, member_id):
        #Inserting persons  image
        if 'female' in img_path : 
            pic2 = slide.shapes.add_picture(img_path, \
        left=Cm(0.20), top=Cm(0.20), width=Cm(5.40), height=Cm(5.20))

        else : pic2 = slide.shapes.add_picture(img_path, \
        left=Cm(0.33), top=Cm(0.27), width=Cm(5.20), height=Cm(5.20))
        
        location_icon =  os.getcwd()+'/static_images/'+ 'location.jpg'
        pic = slide.shapes.add_picture(location_icon,\
              left=Cm(5.86), top=Cm(5.07), width=Cm(1.33), height=Cm(1.33))
        cal_image =  os.getcwd()+'/static_images/'+ 'calender.jpg'
        pic3 = slide.shapes.add_picture(cal_image,\
              left=Cm(5.90), top=Cm(3.95), width=Cm(1.12), height=Cm(1.12))
        p_logo = os.getcwd()+'/static_images/'+ 'pmoves.jpg'
        logo = slide.shapes.add_picture(p_logo, \
              left=Inches(11.81), top=Cm(18.20), width=Cm(3.20), height=Cm(0.85))
        return slide, shapes

    def add_summary(self, slide, shapes, summary, profile_url):
        txBox_sum = slide.shapes.add_textbox(left=Cm(0.10), top=Cm(6.47), \
        width=Cm(0.80), height=Cm(1.60))
        tf = txBox_sum.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'S'
        tf.word_wrap = True
        #tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = 'Gadugi'
        txBox_su = slide.shapes.add_textbox(left=Cm(0.50), top=Cm(6.60), \
        width=Inches(2), height=Inches(0.5))
        tf = txBox_su.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = 'ummary'
        tf.word_wrap = True
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
	txBox_ = slide.shapes.add_textbox(left=Cm(0), \
	top=Cm(8), width=Cm(14), height=Cm(5))
	tf = txBox_.text_frame
	tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
	tf.text = summary[0:1700]
	tf.word_wrap = True
	tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
	tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.name = 'Gadugi'
        icon =  os.getcwd()+'/static_images/'+ 'summary.jpg'
        logo = slide.shapes.add_picture(icon, left=Cm(6.2), top=Cm(6.55), width=Cm(0.92), height=Cm(0.92))
        self.add_connector(slide, shapes, Cm(0.27), Cm(7.64), Cm(7.38), Cm(0.00))

    def add_shapes(self, slide, shapes, RECTANGLE, left, top, width, height):
        shape1 = slide.shapes.add_shape(RECTANGLE, left, top,\
        width, height)
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(216, 191, 216)
        return slide, shapes

    def add_title(self, slide, shapes, title, org, location, exp_data, member_id):
        #Adding Title
        txBox = slide.shapes.add_textbox(left=Cm(5.80), top=Cm(0.10), width=Cm(9.60), height=Cm(1.30))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if options.name == 'yes':
            if title : 
                tf.text = title.replace('  ',' ')
            else : 
                tf.text = 'Name N/A' 
                print "no name available for %s" %member_id

        else : tf.text = 'XXXXXX'
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        if len(title) >= 15 : tf.paragraphs[0].font.size = Pt(16 )
        else : tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
        if len(title) >= 25  : txBox2 = slide.shapes.add_textbox(left=Cm(5.80), top=Cm(1.70), width=Cm(9.76), height=Cm(1.77))
        else : txBox2 = slide.shapes.add_textbox(left=Cm(5.80), top=Cm(1.30), width=Cm(9.76), height=Cm(1.77))
        tf = txBox2.text_frame
        headline = ''
        if  exp_data[0][1] and  exp_data[0][0] : headline = exp_data[0][1] + ' at ' + exp_data[0][0]
        if headline : tf.text = headline[0:100]
        else :  tf.text  = 'Headline N/A'
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.name = 'Gadugi'
        tf.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
        txBox3 = slide.shapes.add_textbox(left=Cm(7), top=Cm(4.20), width=Inches(2.00), height=Cm(1))
        tf = txBox3.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        exp_dur = self.exp_duration(slide, shapes, member_id)
        tf.text = exp_dur
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
        tf.paragraphs[0].font.name = 'Gadugi'
        txBox4 = slide.shapes.add_textbox(left=Cm(7), top=Cm(5.40), width=Inches(2.7), height=Cm(1))
        tf = txBox4.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.text = location
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
        tf.paragraphs[0].font.name = 'Gadugi'
        return slide, shapes

    def exp_duration(self, slide, shapes, member_id):
        self.cur.execute(self.select_dur_qry % member_id)
        exp_dur = self.cur.fetchall()
        exp_dur = exp_dur[0]
        exp_list = []
        if exp_dur:
            for exp in exp_dur:
               if 'years' in exp:
                   dura = "".join(re.findall('\d+', exp)[0])
                   exp_list.append(int(dura))
               elif exp == None : exp_duration = 'Exp N/A'
               else: exp_duration = 'Below 2 years'
            exp_duration = str(sum(exp_list))+'+ Years'
            if '0' in exp_duration: exp_duration = 'Exp N/A'
            print exp_duration
            return exp_duration
    
    def copyrights(self, slide, shapes, page_num):
        page_num_txtbox = slide.shapes.add_textbox(left=Cm(0.30), top=Cm(17.40), width=Cm(1.50), height=Cm(1.89))
        tf = page_num_txtbox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(page_num)
        font = run.font
        font.name = 'Gadugi'
        font.size = Pt(12)
        font.color.rgb = RGBColor(105, 105, 105)
        txtBox = slide.shapes.add_textbox(left=Cm(2.9), top=Inches(6.9), width=Cm(15), height=Cm(1.2))
        tf = txtBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
	run.text = u'\xa9 2017 Positive Moves Consulting | Web: www.positivemoves.com'
        font = run.font
        font.name = 'Gadugi'
        font.size = Pt(12)
        font.color.rgb = RGBColor(105, 105, 105)
 

    def __del__(self):
        self.con.close()
        self.cur.close()

if __name__ == '__main__':
    parser = optparse.OptionParser()
    parser.add_option('-m', '--memberid', default='no', help='member_id, one or many separated by commas')
    parser.add_option('-f', '--filename', default='linkedin_member.pptx', help='filename, give any filename')
    parser.add_option('-i', '--yes', default='no', help='high resolution image /normal image')
    parser.add_option('-t', '--silhoutte', default='no', help='if proper silhoutte needed, type -t yes')
    parser.add_option('-s', '--summary', default='no', help='type -s yes ,if summary is needed ')
    parser.add_option('-e', '--education', default='no', help='type -e yes ,if  education is needed')
    parser.add_option('-p', '--experiences', default='no', help='type -p yes , if experiences if needed')
    parser.add_option('-l', '--languages', default='no', help=' give -l yes, if languages is needed')
    parser.add_option('-n', '--name', default='no', help='type -n yes, if name is needed ')
    (options, args) = parser.parse_args()
    Login(options)

