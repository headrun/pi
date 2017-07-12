import re
import optparse
import MySQLdb
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor
import urllib
import os
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from ppt_db_constants import *

class Login(object):

    def __init__(self, options):
        self.con = MySQLdb.connect(db=db_name,
                      user=db_user,
                      passwd=db_passwd,
                      charset="utf8",
                      host=db_host,
                      use_unicode=True)
        self.cur = self.con.cursor()
        self.x = 1.8 
        self.prs = Presentation()
        self.prs.slide_width = Inches(14.34)
        self.prs.slide_height = Inches(8.5)
        self.select_qry = select_qry
        self.select_qry1 = select_qry1
        self.select_qry2 = select_qry2
        self.select_qry3 = select_qry3
        self.select_qry4 = select_qry4
        self.main()

    def main(self):
        count = 0
        member = options.memberid.split(',')
        for member_id in member:
                self.cur.execute(self.select_qry % member_id)
		rows = self.cur.fetchall()
                if not rows: 
                    print "Invalid member_id %s .Please try by giving valid member_id " % member_id
                    continue
                else:
                    rows = rows[0]
		    name, sk, summary, headline, location = rows
		    self.cur.execute(self.select_qry2 % member_id)
		    profile_image = self.cur.fetchall()
		    try: 
                        path = profile_image[0][1]
                        profile_image = profile_image[0][0]
		    except: profile_image = ''
		    self.get_ppt(rows, profile_image, sk, count, member_id, path)
                    count = count+1

    def get_ppt(self, records, profile_image, sk, count, member_id, path):
        self.cur.execute(self.select_qry1 % sk)
        exp_data = self.cur.fetchall()
        #creating multiple slides
	title_only_slide_layout = self.prs.slide_layouts[0]
	slide = self.prs.slides.add_slide(title_only_slide_layout)
        #slide width
	shapes = slide.shapes
        #Adding title 
        self.add_title(slide, shapes, records[0])
	table = shapes.add_table(rows=4, cols=2, \
        left=Inches(2.4), top=Inches(1), width=Inches(2.0),\
        height=Inches(1.0)).table
	# column widths
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(7.9)
        #Adding headers
        table.cell(0, 0).text = 'Previous Organisation'
        table.cell(0, 1).text = exp_data[0][0]
        table.rows[0].height = Cm(0.2)
        table.rows[1].height = Cm(0.5)
        table.cell(1, 0).text = 'Title/Designation'
	table.cell(1, 1).text = records[3]
        table.rows[2].height = Cm(0.5)
	table.cell(2, 0).text = 'Current Location'
	table.cell(2, 1).text = records[4]
        table.rows[3].height = Cm(0.5)
        table.cell(3, 0).text = 'Work Experience'
        #Borders color
        self.calculate_exp(table, sk)
        #Creating second table 
        table2 = shapes.add_table(rows=4, cols=4, left=Inches(0.0), \
        top=Inches(2.25), width=Inches(6.0), height=Inches(0.5)).table
        table2.columns[0].width = Inches(2.4)
        table2.rows[0].height = Cm(0.5)
        table2.rows[1].height = Cm(0.5)
        table2.rows[2].height = Cm(0.5)
        table2.rows[3].height = Cm(0.5)
        table2.columns[1].width = Inches(7)
        table2.columns[2].width = Inches(2.4)
        table2.columns[3].width = Inches(2.5)
        table2.cell(0, 0).text = '\n\n\n'+ 'Professional experiences'
        for ind, colu in enumerate(exp_data):
                row = 0 
		for indk, i in enumerate(exp_data):
                    if ind == 0 :  
                        paragraphsi_ = [i[0],i[1]]
                        table2.cell(row, 1).text = paragraphsi_[0]
                        table2.cell(row, 1).text_frame.paragraphs[0].runs[0].font.bold = True
                        for para_str in paragraphsi_[1:]:
                                p = table2.cell(row, 1).text_frame.add_paragraph()
                                p.text = para_str
				if indk == 0:
					table2.cell(row, 1).text_frame.paragraphs[1].runs[0].font.bold = False
				else:
	                                p.font.bold = False
                        row = row+1
                    if ind == 1:
                        table2.cell(row, 2).text = i[2]+'\n'+i[3]
                        table2.cell(row, 2).text_frame.paragraphs[0].font.size = Pt(12)
                        row += 1
                    if ind == 2 :
                        table2.cell(row, 3).text = i[4]
                        table2.cell(row, 3).text_frame.paragraphs[0].font.size = Pt(12)
                        row += 1
        self.cur.execute(self.select_qry3 % sk) 
        edu_data = self.cur.fetchall()
        table3 = shapes.add_table(rows=2, cols=3, left=Inches(0.0), \
        top=Inches(4.4), width=Inches(6.0), height=Inches(0.7)).table
        if len(edu_data) < 2:
            table4 = shapes.add_table(rows=1, cols=2, left=Inches(0.0), \
            top=Inches(5.2), width=Inches(6.0), height=Inches(0.5)).table
        else:
            table4 = shapes.add_table(rows=1, cols=2, left=Inches(0.0), \
        top=Inches(5.5), width=Inches(6.0), height=Inches(0.5)).table
        table4.columns[0].width = Inches(2.4)
        table4.rows[0].height = Inches(2.2)
        table4.columns[1].width = Inches(11.9)
        table3.columns[0].width = Inches(2.4)
        table3.rows[0].height = Cm(0.7)
        table3.columns[1].width = Inches(7)
        table3.rows[1].height = Cm(0.7)
        table3.columns[2].width = Inches(4.9)
        table3.cell(0, 0).text = '\n' + 'Education'
        table3.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        table4.cell(0, 0).text = '\n' 'Comments'
        table3.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        #Displaying education data
        for ind, colu in enumerate(edu_data):
                ans = ''
                row = 0
                for inded, i in enumerate(edu_data):
                        if ind == 0: 
                            paragraphsi_ = [i[0],i[1]]
                            table3.cell(row, 1).text = paragraphsi_[0]
                            table3.cell(row, 1).text_frame.paragraphs[0].runs[0].font.bold = True
                            for para_str in paragraphsi_[1:]:
                                p = table3.cell(row, 1).text_frame.add_paragraph()
                                p.text  = para_str
				if inded == 0:
					table3.cell(row, 1).text_frame.paragraphs[1].runs[0].font.bold = False
				else:
	                                p.font.bold = False
                            row = row+1
                        if ind == 1: 
                            table3.cell(row, 2).text = i[2]+'-'+i[3].strip('-')
                            table3.cell(row, 2).text_frame.paragraphs[0].font.size = Pt(12)
                            row += 1
        table4.cell(0, 1).text = records[2][0:1500].encode('utf8')
        tables = [table, table2, table3]
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                   for paragraph in cell.text_frame.paragraphs:
                       for run in paragraph.runs:
                           run.font.size = Pt(13)
                           run.font.color.rgb = RGBColor(0, 0, 0)
        for row in table4.rows:
                for cell in row.cells:
                   for paragraph in cell.text_frame.paragraphs:
                       for run in paragraph.runs:
                           run.font.size = Pt(13)
                           run.font.color.rgb = RGBColor(0, 0, 0)
                           run.font.bold = False
        self.mergeCellsVertically(table2, start_row_idx=0, end_row_idx=3, col_idx=0)
        self.mergeCellsVertically(table3, start_row_idx=0, end_row_idx=1, col_idx=0)
        othr_list_ = [table2.cell(0, 1),table3.cell(0, 1),table2.cell(0, 2),table3.cell(0, 2),table2.cell(0,3),\
        table.cell(0, 0),table4.cell(0,1)]
        for row in othr_list_ :
            row.fill.solid()
            row.fill.fore_color.rgb = RGBColor(230,230,250)
        list_ = [table2.cell(0, 0),table3.cell(0, 0),table4.cell(0, 0)]
        for row in list_ :
            row.fill.solid()
            row.fill.fore_color.rgb = RGBColor(169, 169, 169)
        try :table3.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        except: "no cell"
        try : table4.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True 
        except : print "no cell"
        try :table3.cell(0, 2).text_frame.paragraphs[0].runs[0].font.bold = False
        except : print "no cell"
        try : table2.cell(0, 2).text_frame.paragraphs[0].runs[0].font.bold = False
        except : print "No cell"
        try :table2.cell(0, 3).text_frame.paragraphs[0].runs[0].font.bold = False
        except : print "No cell"
        #Filling colors to cells
        #Downloading image to current working path and save the image with member_id
        img_path = self.download_image(profile_image, member_id, path)
        #Displaying image
        self.insert_image(slide, shapes, img_path, member_id)
        self.copyrights(slide, shapes)
        self.prs.save(options.filename)
        print "Created ppt succesfully for '%s' with name '%s'" % (member_id, records[0].encode('utf-8'))

    def download_image(self, profile_image, member_id, path):
        image = urllib.URLopener()
        pattern = ''
        if '/mpr/mpr/shrink' in profile_image: 
            pattern = "".join(re.findall('mpr/mpr/shrink_\d+_\d+', profile_image))
        if options.yes == 'yes': 
                if pattern: profile_image = profile_image.replace(pattern, 'media')
                else: profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        if options.yes == 'no' and path:
                img_path = path
        if options.yes == 'no' and not path:
                profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        return img_path
        
    def insert_image(self, slide, shapes, img_path, member_id):
        #Inserting persons  image
        pic = slide.shapes.add_picture(img_path,\
              left=Inches(0), top=Inches(1), width=Cm(6.04), height=Cm(4.24)) 
        logo = 'https://positivemoves.com/wp-content/uploads/2017/03/Pmoves-1.png'
        image = urllib.URLopener()
        image_name = image.retrieve(logo, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        pic2 = slide.shapes.add_picture(img_path, \
              left=Inches(12.5), top=Cm(0.80), width=Inches(1.3), height=Cm(1))
        return slide, shapes

    def add_title(self, slide, shapes, title):
        #Adding Title
        txBox = slide.shapes.add_textbox(left=Cm(2.9), top=Cm(0.13), width=Cm(15), height=Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Profile :" + ' ' + title
        font = run.font
        font.bold = True
        font.size = Pt(24)
        line = shapes.add_shape(MSO_CONNECTOR.STRAIGHT,  Cm(0.00), Cm(2.32), Cm(27.90), Cm(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(255, 99, 71) 
        return slide, shapes

    def calculate_exp(self, table, sk):
        exp_duration = 'Below 2 years'
        self.cur.execute(self.select_qry4 % sk)
        exp_dur = self.cur.fetchall()
        exp_list = []
        if exp_dur:
            for exp in exp_dur:
               if 'years' in exp[0]:
                   dura = re.findall('\d+', exp[0])[0]
                   exp_list.append(int(dura))
               else: exp_duration = 'Below 2 years'
            exp_duration = str(sum(exp_list))+'+ Years'
            if '0' in exp_duration: exp_duration = 'Experience Not available'
            table.cell(3, 1).text = exp_duration
            return table

    def merge_cells(self, row_idx, start_col_idx, end_col_idx, table2): 
        col_count = end_col_idx - start_col_idx + 1
        row_cells = [c for c in table2.rows[row_idx].cells][start_col_idx:end_col_idx]
        row_cells[0]._tc.set('gridSpan', str(col_count))
        for c in row_cells[1:]:
            c._tc.set('hMerge', '2')

    def mergeCellsVertically(self, table, start_row_idx, end_row_idx, col_idx):
        row_count = end_row_idx - start_row_idx + 1
        column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]
        column_cells[0]._tc.set('rowSpan', str(row_count))
        for c in column_cells[1:]:
            c._tc.set('vMerge', '1')

    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return parent 
    
    def copyrights(self, slide, shapes):
        txtBox = slide.shapes.add_textbox(left=Cm(2.9), top=Inches(7.9), width=Cm(15), height=Cm(1.2))
        tf = txtBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
	run.text = u'\xa9 2017 Positive Moves Consulting | Web: www.positivemoves.com'
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(13.5)
        font.color.rgb = RGBColor(0, 0, 0)
 
    def border_color(self, cell, border_color="000000", border_width='0.02'):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        lnB = self.SubElement(tcPr, 'a:lnB', w='0.02', cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.SubElement(lnB, 'a:gray')
        srgbClr = self.SubElement(solidFill, 'a:white', val='000000')
        prstDash = self.SubElement(lnB, 'a:prstDash', val='solid')
        round_ = self.SubElement(lnB, 'a:round')
        headEnd = self.SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = self.SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')

    def __del__(self):
        self.con.close()
        self.cur.close()

if __name__ == '__main__':
    parser = optparse.OptionParser()
    parser.add_option('-m', '--memberid', default=None, help='member_id, one or many separated by commas')
    parser.add_option('-f', '--filename', default='linkedin_member.pptx', help='filename, give any filename')
    parser.add_option('-i', '--yes', default='no', help='high resolution image /normal image')
    (options, args) = parser.parse_args()
    Login(options)

