import re
import optparse
import MySQLdb
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor
import urllib
import os
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from ppt_db_constants import *


class Login(object):

    def __init__(self, options):
        self.con = MySQLdb.connect(db=db_name,
                      user=db_user,
                      passwd=db_passwd,
                      charset="utf8",
                      host=db_host,
                      use_unicode=True)
        self.cur = self.con.cursor()
        self.prs = Presentation()
        self.prs.slide_width = Inches(14.34)
        self.prs.slide_height = Inches(8.5)
        self.select_qry = select_qry
        self.select_qry1 = select_qry1
        self.select_qry2 = select_qry2
        self.select_qry3 = select_qry3
        self.select_qry4 = select_qry4
        self.main()

    def main(self):
        count = 0
        member = options.memberid.split(',')
        for member_id in member:
                self.cur.execute(self.select_qry % member_id)
		rows = self.cur.fetchall()
                if not rows: 
                    print "Invalid member_id %s .Please try by giving valid member_id " % member_id
                    continue
                else:
                    rows = rows[0]
		    name, sk, summary, headline, location = rows
		    self.cur.execute(self.select_qry2 % member_id)
		    profile_image = self.cur.fetchall()
		    try: 
                        path = profile_image[0][1]
                        profile_image = profile_image[0][0]
		    except: profile_image = ''
		    self.get_ppt(rows, profile_image, sk, count, member_id, path)
                    count = count+1


    def get_ppt(self, records, profile_image, sk, count, member_id, path):
        self.cur.execute(self.select_qry1 % sk)
        exp_data = self.cur.fetchall()
        #creating multiple slides
	title_only_slide_layout = self.prs.slide_layouts[0]
	slide = self.prs.slides.add_slide(title_only_slide_layout)
        #slide width
	shapes = slide.shapes
        #Adding title 
        self.add_title(slide, shapes, records[0])
	table = shapes.add_table(rows=4, cols=2, \
        left=Inches(2.4), top=Inches(1), width=Inches(2.0),\
        height=Inches(1.0)).table
	# column widths
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(7.9)
        #Adding headers
        table.cell(0, 0).text = 'Previous Organisation'
        table.cell(0, 1).text = exp_data[0][1]
        table.rows[0].height = Cm(0.5)
        table.rows[1].height = Cm(0.5)
        table.cell(1, 0).text = 'Title/Designation'
	table.cell(1, 1).text = records[3]
        table.rows[2].height = Cm(0.5)
	table.cell(2, 0).text = 'Current Location'
	table.cell(2, 1).text = records[4]
        table.rows[3].height = Cm(0.5)
        table.cell(3, 0).text = 'Work Experience'
        #Borders color
        for row in table.rows:
                for cell in row.cells:
                    self.border_color(cell)
        self.calculate_exp(table, sk)
        #Creating second table 
        table2 = shapes.add_table(rows=3, cols=4, left=Inches(0.0), \
        top=Inches(2.1), width=Inches(6.0), height=Inches(0.5)).table
        table2.columns[0].width = Inches(2.4)
        table2.rows[0].height = Inches(2)
        table2.columns[1].width = Inches(5.3)
        table2.rows[1].height = Inches(1.2)
        table2.columns[2].width = Inches(3.8)
        table2.columns[3].width = Inches(2.8)
        table2.rows[2].height = Inches(2)
        table2.cell(0, 0).text = '\n\n\n'+ 'Professional experiences'
        table2.cell(1, 0).text = 'Education'
        table2.cell(2, 0).text = 'comments'
        #Filling colors to cells
        tables = [table, table2]
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(70, 130, 180)
        row_idx, start_col_idx, end_col_idx = 2, 1, 3
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table2)
        row_idx, start_col_idx, end_col_idx = 1, 1, 2
        self.merge_cells(row_idx, start_col_idx, end_col_idx, table2)
        #Displaying experiences data
	for ind, colu in enumerate(exp_data):
		ans = ''
		for i in exp_data:
			if ind == 0: ans += '\n'.join(list(i[0:2]))+'\n\n'
			if ind == 1: ans += '\n'.join(list(i[2:4]))+'\n\n'
                        if ind == 2: ans += "".join(list(i[4]))+'\n\n\n'
		try: 
                    table2.cell(0, ind+1).text = ans.strip('\n')
                except: break
        #Displaying education data
        self.cur.execute(self.select_qry3 % sk)
        edu_data = self.cur.fetchall()
        for ind, colu in enumerate(edu_data):
                ans = ''
                for i in edu_data:
                        if ind == 0: ans += '\n'.join(list(i[0:2]))+'\n\n'
                        if ind == 1: ans += '-'.join(list(i[2:4]))+'\n\n'
                try:
                    if ind == 1: table2.cell(1, ind+2).text = str(ans).strip('\n').strip('-')
                    if ind == 0: table2.cell(1, ind+1).text = ans.strip('\n')
                except: break
        table2.cell(2, 1).text = records[2][0:1500].encode('utf8')
        #Adding font color 
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                   for paragraph in cell.text_frame.paragraphs:
                       for run in paragraph.runs:
                           run.font.size = Pt(11)
                           run.font.color.rgb = RGBColor(0, 0, 0)
                           run.font.bold = False
        list_ = [table2.cell(0, 0), table2.cell(1, 0), table2.cell(2, 0)]
        for row in list_:
            row.fill.solid()
            row.fill.fore_color.rgb = RGBColor(169, 169, 169)
        for cell in list_:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
        for row in table2.rows:
                for cell in row.cells:
                    self.border_color(cell)
        #Downloading image  to current working path and save the image with member_id
        img_path = self.download_image(profile_image, member_id, path)
        #Displaying image
        self.insert_image(slide, shapes, img_path, member_id)
        self.copyrights(slide, shapes)
        self.prs.save(options.filename)
        print "Created ppt succesfully for '%s' with name '%s'" % (member_id, records[0].encode('utf-8'))

    def download_image(self, profile_image, member_id, path):
        image = urllib.URLopener()
        pattern = ''
        if '/mpr/mpr/shrink' in profile_image: 
            pattern = "".join(re.findall('mpr/mpr/shrink_\d+_\d+', profile_image))
        if options.yes == 'yes': 
                if pattern: profile_image = profile_image.replace(pattern, 'media')
                else: profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        if options.yes == 'no' and path:
                img_path = path
        if options.yes == 'no' and not path:
                profile_image = 'http://bento.cdn.pbs.org/hostedbento-prod/filer_public/_bento_media/img/no-image-available.jpg'
                image_name = image.retrieve(profile_image, '%s.jpg'% member_id)
                img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        return img_path
        

    def insert_image(self, slide, shapes, img_path, member_id):
        #Inserting persons  image
        pic = slide.shapes.add_picture(img_path,\
              left=Inches(0), top=Inches(1), width=Cm(6.04), height=Cm(4.24)) 
        logo = 'https://positivemoves.com/wp-content/uploads/2017/03/Pmoves-1.png'
        image = urllib.URLopener()
        image_name = image.retrieve(logo, 'logo.jpg')
        img_path = os.path.dirname(os.path.abspath(image_name[0]))+ '/' +image_name[0]
        pic2 = slide.shapes.add_picture(img_path, \
              left=Inches(12.5), top=Cm(0.80), width=Inches(1.3), height=Cm(1))
        return slide, shapes

    def add_title(self, slide, shapes, title):
        #Adding Title
        txBox = slide.shapes.add_textbox(left=Cm(2.9), top=Cm(0.13), width=Cm(15), height=Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Profile :" + ' ' + title
        font = run.font
        font.bold = True
        font.size = Pt(24)
        line = shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER, Cm(0), Cm(2.32), Cm(30.80), Cm(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(255, 99, 71) 
        return slide, shapes

    def calculate_exp(self, table, sk):
        exp_duration = 'Below 2 years'
        self.cur.execute(self.select_qry4 % sk)
        exp_dur = self.cur.fetchall()
        exp_list = []
        if exp_dur:
            for exp in exp_dur:
               if 'years' in exp[0]:
                   dura = re.findall('\d+', exp[0])[0]
                   exp_list.append(int(dura))
               else: exp_duration = 'Below 2 years'
            exp_duration = str(sum(exp_list))+'+ Years'
            if '0' in exp_duration: exp_duration = 'Experience Not available'
            table.cell(3, 1).text = exp_duration
            return table

    def merge_cells(self, row_idx, start_col_idx, end_col_idx, table2): 
        col_count = end_col_idx - start_col_idx + 1
        row_cells = [c for c in table2.rows[row_idx].cells][start_col_idx:end_col_idx]
        row_cells[0]._tc.set('gridSpan', str(col_count))
        for c in row_cells[1:]:
            c._tc.set('hMerge', '2')


    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return parent 
    
    def copyrights(self, slide, shapes):
        txtBox = slide.shapes.add_textbox(left=Cm(2.9), top=Inches(7.9), width=Cm(15), height=Cm(1.2))
        tf = txtBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
	run.text = u'\xa9 2017 Positive Moves Consulting | Web: www.positivemoves.com'
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(13.5)
        font.color.rgb = RGBColor(0, 0, 0)
 
    def border_color(self, cell, border_color="000000", border_width='0.02'):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        lnB = self.SubElement(tcPr, 'a:lnB', w='0.02', cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.SubElement(lnB, 'a:gray')
        srgbClr = self.SubElement(solidFill, 'a:white', val='000000')
        prstDash = self.SubElement(lnB, 'a:prstDash', val='solid')
        round_ = self.SubElement(lnB, 'a:round')
        headEnd = self.SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = self.SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')


    def __del__(self):
        self.con.close()
        self.cur.close()

if __name__ == '__main__':
    parser = optparse.OptionParser()
    parser.add_option('-m', '--memberid', default=None, help='member_id, one or many separated by commas')
    parser.add_option('-f', '--filename', default='linkedin_member.pptx', help='filename, give any filename')
    parser.add_option('-i', '--yes', default='no', help='high resolution image /normal image')
    (options, args) = parser.parse_args()
    Login(options)

