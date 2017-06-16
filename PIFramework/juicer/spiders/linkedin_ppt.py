import os
import optparse
import MySQLdb
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

class Login(object):

    def __init__(self,options):

        self.con = MySQLdb.connect(db='FACEBOOK',
                      user='root',
                      passwd='root',
                      charset="utf8",
                      host='localhost',
                      use_unicode=True)

        self.cur = self.con.cursor()
        self.prs = Presentation()
        self.sample_text = "test.pptx"
        self.select_qry = 'select name,first_name,last_name,summary,headline,location from linkedin_meta where member_id = "%s" '
        self.main()
        

    def main(self):
        member_id = int(options.memberid)
        self.cur.execute(self.select_qry%member_id)
        rows = self.cur.fetchall()
        rows = rows[0]
        name,first_name,last_name,summary,headline,location = rows
        self.get_ppt(rows)

    def get_ppt(self,records):
	title_only_slide_layout = self.prs.slide_layouts[0]
	slide = self.prs.slides.add_slide(title_only_slide_layout)
	shapes = slide.shapes
	shapes.title.text = records[0]
	rows = cols = 4
	left = Inches(0.0)
        top = Inches(0)
	width = Inches(5.0)
	height = Inches(0.5)
       
	table = shapes.add_table(rows, cols, left, top, width, height).table
	# set column widths
	table.columns[0].width = Inches(1.8)
	table.columns[1].width = Inches(9.0)
        	
        table.cell(0, 0).text  = 'Title/designation'
	table.cell(0, 1).text  = records[4]

	table.cell(1, 0).text = 'Current Location'
	table.cell(1, 1).text = records[5]

        table.cell(2,0).text = 'Experience'
        table.cell(2,1).text = records[3]
        self.prs.save('test.pptx')
        img_path = '/root/Linkedin/Linkedin/spiders/images/full/d43302b6dbbc9c21b248bbb74c184afa17b6d74c.jpg'
        

    def __del__(self):
        self.con.close()
        self.cur.close()

   

if __name__ == '__main__':
    parser = optparse.OptionParser()
    parser.add_option('-m', '--memberid', default=None, help='member_id')
    (options, args) = parser.parse_args()
    Login(options)

