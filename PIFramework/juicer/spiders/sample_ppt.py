from pptx import Presentation
from pptx.util import Inches

class FirstPpt(object):

    def __init__(self):
        self.prs = Presentation()
        self.sample_text = 'sample.pptx'

    def first_slide(self):
	title_slide_layout = self.prs.slide_layouts[0]
	slide = self.prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = "Hello, PI Team "
	subtitle.text = "This is our first slide by pptx module"
	self.prs.save(self.sample_text)

    def second_slide(self):
	bullet_slide_layout = self.prs.slide_layouts[1]
	slide = self.prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = 'Second slide by PPTX'
	tf = body_shape.text_frame
	tf.text ='Hello Aravind/kiranmayi'
	p = tf.add_paragraph()
	p.text = 'This is second Slide by PPTX module.'
	p.level = 1
	p = tf.add_paragraph()
	p.text = 'I have read this document and tried this example'
	p.level = 2
	self.prs.save(self.sample_text)

    def third_slide(self):
        
        img_path = '/root/Linkedin/Linkedin/spiders/images/full/d43302b6dbbc9c21b248bbb74c184afa17b6d74c.jpg'
	blank_slide_layout = self.prs.slide_layouts[3]
	slide = self.prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Adding picture to this slide'
	left = top = Inches(2)
	height = Inches(3.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
	self.prs.save(self.sample_text)


    def main(self):

        first_slide = self.first_slide()
        Add_bullet = self.second_slide()
        Add_picture = self.third_slide()

if __name__ == "__main__":
     FirstPpt().main()





